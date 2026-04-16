"""
Üzleti Központ Blueprint
ERP kiegészítő modul: kiadások, munkás kifizetések, leadek, feladatok, teljesítési igazolások
"""
from flask import Blueprint, render_template, request, flash, redirect, url_for, session, jsonify
from datetime import date, datetime, timedelta
import logging
from app.extensions import db
from app.utils.decorators import handle_database_errors, log_function_call
from app.utils.validators import sanitize_input

business_bp = Blueprint('business', __name__)
logger = logging.getLogger(__name__)


def require_login():
    if not session.get('logged_in'):
        flash('Kérjük jelentkezzen be!', 'warning')
        return redirect(url_for('auth.login'))
    return None


def get_tenant_id():
    return session.get('current_tenant_id') or 1


def get_workers():
    from app.models.user import User
    from app.models.tenant_membership import TenantMembership
    from sqlalchemy import and_
    tenant_id = get_tenant_id()
    try:
        memberships = db.session.execute(
            db.select(TenantMembership).where(
                and_(TenantMembership.tenant_id == tenant_id, TenantMembership.status == 'active')
            )
        ).scalars().all()
        user_ids = [m.user_id for m in memberships]
        if not user_ids:
            return []
        return db.session.execute(
            db.select(User).where(User.user_id.in_(user_ids))
        ).scalars().all()
    except Exception as e:
        logger.error(f"get_workers error: {e}")
        return []


# ─────────────────────────────────────────────────────────────────
# MAIN HUB
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub')
@handle_database_errors
def hub():
    r = require_login()
    if r: return r

    tenant_id = get_tenant_id()
    today = date.today()

    try:
        from app.models.expense import Expense
        from app.models.lead import Lead
        from app.models.task import Task
        from app.models.worker_payment import WorkerPayment, PerformanceConfirmation
        from app.models.job import Job
        from app.models.user import User
        from sqlalchemy import and_, func, case

        # ── HAVI ÖSSZESÍTŐK ────────────────────────────────────────
        total_expenses = db.session.execute(
            db.select(func.coalesce(func.sum(Expense.amount), 0))
            .where(and_(Expense.tenant_id == tenant_id,
                        func.extract('month', Expense.expense_date) == today.month,
                        func.extract('year', Expense.expense_date) == today.year))
        ).scalar() or 0

        total_payouts = db.session.execute(
            db.select(func.coalesce(func.sum(WorkerPayment.amount), 0))
            .where(and_(WorkerPayment.tenant_id == tenant_id,
                        func.extract('month', WorkerPayment.payment_date) == today.month,
                        func.extract('year', WorkerPayment.payment_date) == today.year))
        ).scalar() or 0

        open_leads = db.session.execute(
            db.select(func.count()).select_from(Lead)
            .where(and_(Lead.tenant_id == tenant_id, Lead.stage.notin_(['Megnyert', 'Elveszett'])))
        ).scalar() or 0

        open_tasks = db.session.execute(
            db.select(func.count()).select_from(Task)
            .where(and_(Task.tenant_id == tenant_id, Task.done == False))
        ).scalar() or 0

        monthly_revenue = db.session.execute(
            db.select(func.coalesce(func.sum(Job.total_cost), 0))
            .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                        func.extract('month', Job.job_date) == today.month,
                        func.extract('year', Job.job_date) == today.year))
        ).scalar() or 0

        # ── PROFIT RÉSZLETEZÉS ─────────────────────────────────────
        # Anyagköltség (alkatrészek)
        try:
            from app.models.job import JobPart
            material_cost = db.session.execute(
                db.select(func.coalesce(
                    func.sum(JobPart.quantity * JobPart.unit_price), 0))
                .join(Job, JobPart.job_id == Job.job_id)
                .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                            func.extract('month', Job.job_date) == today.month,
                            func.extract('year', Job.job_date) == today.year))
            ).scalar() or 0
        except Exception:
            material_cost = 0

        labor_cost = float(total_payouts)
        net_profit = float(monthly_revenue) - float(total_expenses) - labor_cost
        profit_pct = round((net_profit / float(monthly_revenue) * 100), 1) if float(monthly_revenue) > 0 else 0

        # ── MAI TEENDŐK ────────────────────────────────────────────
        today_jobs = db.session.execute(
            db.select(func.count()).select_from(Job)
            .where(and_(Job.tenant_id == tenant_id,
                        Job.job_date == today,
                        Job.completed == False))
        ).scalar() or 0

        overdue_tasks = db.session.execute(
            db.select(Task)
            .where(and_(Task.tenant_id == tenant_id,
                        Task.done == False,
                        Task.deadline != None,
                        Task.deadline <= today))
            .order_by(Task.deadline)
            .limit(5)
        ).scalars().all()

        today_tasks = db.session.execute(
            db.select(Task)
            .where(and_(Task.tenant_id == tenant_id,
                        Task.done == False,
                        Task.deadline == today))
        ).scalars().all()

        # Lejárt/fizetetlen számlák (jobs where completed but not invoiced / old)
        try:
            unpaid_jobs = db.session.execute(
                db.select(Job)
                .where(and_(Job.tenant_id == tenant_id,
                            Job.completed == True,
                            Job.job_date <= today - timedelta(days=30)))
                .order_by(Job.job_date.asc())
                .limit(5)
            ).scalars().all()
        except Exception:
            unpaid_jobs = []

        # Kifizetendő munkások (van-e PerformanceConfirmation 'Beadva' státuszban)
        try:
            pending_confirmations = db.session.execute(
                db.select(func.count()).select_from(PerformanceConfirmation)
                .where(and_(PerformanceConfirmation.tenant_id == tenant_id,
                            PerformanceConfirmation.status == 'Beadva'))
            ).scalar() or 0
        except Exception:
            pending_confirmations = 0

        # ── CASHFLOW ADATOK (30 nap) — egyetlen lekérdezéssel ────────
        cashflow_days = []
        date_from = today - timedelta(days=29)

        # Bevétel naponta
        rev_rows = db.session.execute(
            db.select(Job.job_date, func.coalesce(func.sum(Job.total_cost), 0).label('rev'))
            .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                        Job.job_date >= date_from, Job.job_date <= today))
            .group_by(Job.job_date)
        ).all()
        rev_map = {r.job_date: float(r.rev) for r in rev_rows}

        # Kiadás naponta
        exp_rows = db.session.execute(
            db.select(Expense.expense_date, func.coalesce(func.sum(Expense.amount), 0).label('exp'))
            .where(and_(Expense.tenant_id == tenant_id,
                        Expense.expense_date >= date_from, Expense.expense_date <= today))
            .group_by(Expense.expense_date)
        ).all()
        exp_map = {r.expense_date: float(r.exp) for r in exp_rows}

        for i in range(29, -1, -1):
            d = today - timedelta(days=i)
            rev = rev_map.get(d, 0.0)
            exp = exp_map.get(d, 0.0)
            cashflow_days.append({
                'date': d.strftime('%m.%d'),
                'revenue': rev,
                'expense': exp,
                'profit': rev - exp
            })

        # ── TOP MUNKÁSOK — napi/heti/havi bontással ───────────────
        workers = get_workers()
        top_workers = []
        avg_jobs = 0

        # Periódus: napi / heti / havi (default: havi)
        period = request.args.get('period', 'monthly')
        if period == 'daily':
            period_start = today
            period_label = 'Ma'
        elif period == 'weekly':
            period_start = today - timedelta(days=today.weekday())
            period_label = 'Ezen a héten'
        else:
            period = 'monthly'
            period_start = today.replace(day=1)
            period_label = today.strftime('%Y. %B')

        try:
            worker_stats = []
            for w in workers:
                # Munkák száma az adott periódusban
                if period == 'daily':
                    wc = db.session.execute(
                        db.select(func.count()).select_from(Job)
                        .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                                    Job.job_date == today))
                    ).scalar() or 0
                    rev = db.session.execute(
                        db.select(func.coalesce(func.sum(Job.total_cost), 0))
                        .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                                    Job.job_date == today))
                    ).scalar() or 0
                else:
                    wc = db.session.execute(
                        db.select(func.count()).select_from(Job)
                        .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                                    Job.job_date >= period_start))
                    ).scalar() or 0
                    rev = db.session.execute(
                        db.select(func.coalesce(func.sum(Job.total_cost), 0))
                        .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                                    Job.job_date >= period_start))
                    ).scalar() or 0

                # Trend: előző periódus
                if period == 'daily':
                    prev_start = today - timedelta(days=1)
                    prev_end   = today - timedelta(days=1)
                elif period == 'weekly':
                    prev_start = period_start - timedelta(weeks=1)
                    prev_end   = period_start - timedelta(days=1)
                else:
                    prev_month = (today.replace(day=1) - timedelta(days=1))
                    prev_start = prev_month.replace(day=1)
                    prev_end   = prev_month

                prev_wc = db.session.execute(
                    db.select(func.count()).select_from(Job)
                    .where(and_(Job.tenant_id == tenant_id, Job.completed == True,
                                Job.job_date >= prev_start, Job.job_date <= prev_end))
                ).scalar() or 0

                trend = 'up' if wc > prev_wc else ('down' if wc < prev_wc else 'flat')
                trend_pct = round((wc - prev_wc) / prev_wc * 100) if prev_wc > 0 else (100 if wc > 0 else 0)

                worker_stats.append({
                    'name': getattr(w, 'username', 'Ismeretlen'),
                    'job_count': wc,
                    'revenue': float(rev),
                    'trend': trend,
                    'trend_pct': trend_pct,
                    'prev_count': prev_wc,
                })

            worker_stats.sort(key=lambda x: x['job_count'], reverse=True)
            top_workers = worker_stats[:5]
            if worker_stats:
                avg_jobs = sum(s['job_count'] for s in worker_stats) / len(worker_stats)
        except Exception as e:
            logger.warning(f"top_workers error: {e}")

        # ── FIGYELMEZTETÉSEK ───────────────────────────────────────
        warnings = []

        # Hiányzó teljesítési igazolások
        if pending_confirmations > 0:
            warnings.append({
                'type': 'danger',
                'icon': 'ti-file-alert',
                'text': f'{pending_confirmations} teljesítési igazolás jóváhagyásra vár',
                'url': url_for('business.confirmations')
            })

        # Lejárt feladatok
        overdue_count = len(overdue_tasks)
        if overdue_count > 0:
            warnings.append({
                'type': 'warning',
                'icon': 'ti-clock-exclamation',
                'text': f'{overdue_count} lejárt feladat',
                'url': url_for('business.tasks')
            })

        # Negatív profit
        if net_profit < 0:
            warnings.append({
                'type': 'danger',
                'icon': 'ti-trending-down',
                'text': f'Negatív havi profit: {int(net_profit):,} Ft'.replace(',', ' '),
                'url': url_for('business.expenses')
            })

        # Sok lead veszteség
        lost_leads = db.session.execute(
            db.select(func.count()).select_from(Lead)
            .where(and_(Lead.tenant_id == tenant_id, Lead.stage == 'Elveszett',
                        func.extract('month', Lead.created_date) == today.month))
        ).scalar() or 0
        if lost_leads >= 3:
            warnings.append({
                'type': 'warning',
                'icon': 'ti-user-x',
                'text': f'{lost_leads} elveszett lead ebben a hónapban',
                'url': url_for('business.leads')
            })

        # ── LEGUTÓBBI AKTIVITÁSOK ──────────────────────────────────
        recent_activities = []

        # Utolsó 5 kiadás
        recent_expenses = db.session.execute(
            db.select(Expense).where(Expense.tenant_id == tenant_id)
            .order_by(Expense.expense_date.desc()).limit(3)
        ).scalars().all()
        for e in recent_expenses:
            recent_activities.append({
                'icon': 'ti-receipt',
                'color': 'danger',
                'text': f'Kiadás rögzítve: {int(e.amount):,} Ft ({e.category or "egyéb"})'.replace(',', ' '),
                'date': e.expense_date.strftime('%m.%d') if e.expense_date else '',
                'sort_date': e.expense_date or date.min
            })

        # Utolsó 3 befejezett munka
        recent_jobs = db.session.execute(
            db.select(Job).where(and_(Job.tenant_id == tenant_id, Job.completed == True))
            .order_by(Job.job_date.desc()).limit(3)
        ).scalars().all()
        for j in recent_jobs:
            recent_activities.append({
                'icon': 'ti-check',
                'color': 'success',
                'text': f'Munka lezárva: {int(j.total_cost or 0):,} Ft'.replace(',', ' '),
                'date': j.job_date.strftime('%m.%d') if j.job_date else '',
                'sort_date': j.job_date or date.min
            })

        # Utolsó 2 lead
        recent_leads_list = db.session.execute(
            db.select(Lead).where(Lead.tenant_id == tenant_id)
            .order_by(Lead.created_date.desc()).limit(2)
        ).scalars().all()
        for l in recent_leads_list:
            recent_activities.append({
                'icon': 'ti-user-plus',
                'color': 'primary',
                'text': f'Új lead: {l.name} ({l.stage})',
                'date': l.created_date.strftime('%m.%d') if l.created_date else '',
                'sort_date': l.created_date or date.min
            })

        recent_activities.sort(key=lambda x: x['sort_date'], reverse=True)
        recent_activities = recent_activities[:8]

        # ── AJÁNLATKÉSZÍTŐHÖZ: szolgáltatások, alkatrészek, ügyfelek ──
        try:
            from app.models.service import Service
            services_list = db.session.execute(
                db.select(Service).where(Service.tenant_id == tenant_id)
                .order_by(Service.service_name)
            ).scalars().all()
        except Exception:
            services_list = []
        try:
            from app.models.part import Part
            parts_list = db.session.execute(
                db.select(Part).where(Part.tenant_id == tenant_id)
                .order_by(Part.part_name)
            ).scalars().all()
        except Exception:
            parts_list = []
        try:
            quote_jobs = db.session.execute(
                db.select(Job).where(Job.tenant_id == tenant_id)
                .order_by(Job.job_id.desc()).limit(30)
            ).scalars().all()
        except Exception:
            quote_jobs = []

        return render_template('business/hub.html',
            # Havi összesítők
            total_expenses=float(total_expenses),
            total_payouts=float(total_payouts),
            open_leads=open_leads,
            open_tasks=open_tasks,
            monthly_revenue=float(monthly_revenue),
            profit_estimate=net_profit,
            current_month=today.strftime('%Y. %B'),
            # Profit részletezés
            material_cost=float(material_cost),
            labor_cost=labor_cost,
            net_profit=net_profit,
            profit_pct=profit_pct,
            # Mai teendők
            today_jobs=today_jobs,
            overdue_tasks=overdue_tasks,
            today_tasks=today_tasks,
            unpaid_jobs=unpaid_jobs,
            pending_confirmations=pending_confirmations,
            # Cashflow
            cashflow_days=cashflow_days,
            # Top munkások
            top_workers=top_workers,
            avg_jobs=avg_jobs,
            period=period,
            period_label=period_label,
            # Figyelmeztetések
            warnings=warnings,
            # Aktivitások
            recent_activities=recent_activities,
            # Workers (modal-hoz)
            workers=workers,
            today=today,
            # AI ajánlatkészítő
            services_list=services_list,
            parts_list=parts_list,
            recent_jobs=quote_jobs,
        )

    except Exception as e:
        logger.error(f"Business hub error: {e}", exc_info=True)
        flash(f'Dashboard hiba: {e}', 'error')
        return render_template('business/hub.html',
            total_expenses=0, total_payouts=0, open_leads=0, open_tasks=0,
            monthly_revenue=0, profit_estimate=0, current_month=today.strftime('%Y. %B'),
            material_cost=0, labor_cost=0, net_profit=0, profit_pct=0,
            today_jobs=0, overdue_tasks=[], today_tasks=[], unpaid_jobs=[],
            pending_confirmations=0, cashflow_days=[], top_workers=[], avg_jobs=0,
            warnings=[], recent_activities=[], workers=[], today=today)


# ─────────────────────────────────────────────────────────────────
# NAPI ZÁRÁS API
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/daily-close', methods=['POST'])
@handle_database_errors
def daily_close():
    r = require_login()
    if r: return jsonify({'error': 'Nincs bejelentkezve'}), 401

    tenant_id = get_tenant_id()
    today = date.today()
    issues = []
    ok = []

    try:
        from app.models.task import Task
        from app.models.expense import Expense
        from app.models.job import Job
        from sqlalchemy import and_, func

        # Nyitott feladatok ma
        open_tasks_count = db.session.execute(
            db.select(func.count()).select_from(Task)
            .where(and_(Task.tenant_id == tenant_id, Task.done == False,
                        Task.deadline != None, Task.deadline <= today))
        ).scalar() or 0

        if open_tasks_count > 0:
            issues.append(f'⚠️ {open_tasks_count} lejárt feladat maradt nyitva')
        else:
            ok.append('✅ Minden lejárt feladat elvégezve')

        # Mai befejezetlen munkák
        open_jobs = db.session.execute(
            db.select(func.count()).select_from(Job)
            .where(and_(Job.tenant_id == tenant_id, Job.completed == False,
                        Job.job_date == today))
        ).scalar() or 0

        if open_jobs > 0:
            issues.append(f'⚠️ {open_jobs} mai munka nincs lezárva')
        else:
            ok.append('✅ Minden mai munka lezárva')

        # Rögzítetlen kiadások ellenőrzése (ha nincs egyetlen mai kiadás sem, és van munka)
        today_expenses = db.session.execute(
            db.select(func.count()).select_from(Expense)
            .where(and_(Expense.tenant_id == tenant_id, Expense.expense_date == today))
        ).scalar() or 0

        today_completed = db.session.execute(
            db.select(func.count()).select_from(Job)
            .where(and_(Job.tenant_id == tenant_id, Job.completed == True, Job.job_date == today))
        ).scalar() or 0

        if today_completed > 0 and today_expenses == 0:
            issues.append('💡 Tipp: Volt munka ma, de nem lett kiadás rögzítve')
        elif today_expenses > 0:
            ok.append(f'✅ {today_expenses} kiadás rögzítve ma')

        status = 'ok' if len(issues) == 0 else 'warning'
        return jsonify({
            'status': status,
            'issues': issues,
            'ok': ok,
            'date': today.strftime('%Y. %m. %d.')
        })

    except Exception as e:
        return jsonify({'status': 'error', 'issues': [str(e)], 'ok': []}), 500


# ─────────────────────────────────────────────────────────────────
# KIADÁSOK
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/expenses')
@handle_database_errors
def expenses():
    r = require_login()
    if r: return r
    from app.models.expense import Expense
    from sqlalchemy import and_
    tenant_id = get_tenant_id()
    items = db.session.execute(
        db.select(Expense).where(Expense.tenant_id == tenant_id).order_by(Expense.expense_date.desc())
    ).scalars().all()
    workers = get_workers()
    return render_template('business/expenses.html', expenses=items, workers=workers, today=date.today())


@business_bp.route('/business-hub/expenses/add', methods=['POST'])
@handle_database_errors
def add_expense():
    r = require_login()
    if r: return r
    from app.models.expense import Expense
    try:
        e = Expense(
            tenant_id=get_tenant_id(),
            worker_id=request.form.get('worker_id', type=int) or None,
            amount=float(request.form.get('amount', 0)),
            currency=request.form.get('currency', 'HUF'),
            payment_source=sanitize_input(request.form.get('payment_source', '')),
            category=sanitize_input(request.form.get('category', '')),
            notes=sanitize_input(request.form.get('notes', '')),
            expense_date=date.fromisoformat(request.form.get('expense_date', str(date.today()))),
        )
        db.session.add(e)
        db.session.commit()
        flash('Kiadás rögzítve!', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    next_url = request.form.get('next', url_for('business.expenses'))
    return redirect(next_url)


@business_bp.route('/business-hub/expenses/<int:expense_id>/delete', methods=['POST'])
@handle_database_errors
def delete_expense(expense_id):
    r = require_login()
    if r: return r
    from app.models.expense import Expense
    e = db.session.get(Expense, expense_id)
    if e:
        db.session.delete(e)
        db.session.commit()
        flash('Kiadás törölve!', 'success')
    return redirect(url_for('business.expenses'))


# ─────────────────────────────────────────────────────────────────
# MUNKÁS KIFIZETÉSEK
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/worker-payments')
@handle_database_errors
def worker_payments():
    r = require_login()
    if r: return r
    from app.models.worker_payment import WorkerPayment
    from sqlalchemy import and_, func
    tenant_id = get_tenant_id()
    workers = get_workers()

    payment_totals = {}
    for w in workers:
        total = db.session.execute(
            db.select(func.coalesce(func.sum(WorkerPayment.amount), 0))
            .where(and_(WorkerPayment.tenant_id == tenant_id, WorkerPayment.worker_id == w.user_id))
        ).scalar() or 0
        payment_totals[w.user_id] = float(total)

    payments = db.session.execute(
        db.select(WorkerPayment).where(WorkerPayment.tenant_id == tenant_id)
        .order_by(WorkerPayment.payment_date.desc())
    ).scalars().all()

    return render_template('business/worker_payments.html',
        workers=workers, payments=payments, payment_totals=payment_totals, today=date.today())


@business_bp.route('/business-hub/worker-payments/add', methods=['POST'])
@handle_database_errors
def add_worker_payment():
    r = require_login()
    if r: return r
    from app.models.worker_payment import WorkerPayment
    try:
        worker_id = request.form.get('worker_id', type=int)
        if not worker_id:
            flash('Munkás kiválasztása kötelező!', 'error')
            next_url = request.form.get('next', url_for('business.worker_payments'))
            return redirect(next_url)
        p = WorkerPayment(
            tenant_id=get_tenant_id(),
            worker_id=worker_id,
            amount=float(request.form.get('amount', 0)),
            currency=request.form.get('currency', 'HUF'),
            payment_source=sanitize_input(request.form.get('payment_source', 'Készpénz')),
            payment_date=date.fromisoformat(request.form.get('payment_date', str(date.today()))),
            notes=sanitize_input(request.form.get('notes', '')),
        )
        db.session.add(p)
        db.session.commit()
        flash('Kifizetés rögzítve!', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    next_url = request.form.get('next', url_for('business.worker_payments'))
    return redirect(next_url)


@business_bp.route('/business-hub/worker-payments/<int:payment_id>/delete', methods=['POST'])
@handle_database_errors
def delete_worker_payment(payment_id):
    r = require_login()
    if r: return r
    from app.models.worker_payment import WorkerPayment
    p = db.session.get(WorkerPayment, payment_id)
    if p:
        db.session.delete(p)
        db.session.commit()
        flash('Kifizetés törölve!', 'success')
    return redirect(url_for('business.worker_payments'))


# ─────────────────────────────────────────────────────────────────
# LEADEK
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/leads')
@handle_database_errors
def leads():
    r = require_login()
    if r: return r
    from app.models.lead import Lead, LEAD_STAGES
    from sqlalchemy import and_
    tenant_id = get_tenant_id()
    all_leads = db.session.execute(
        db.select(Lead).where(Lead.tenant_id == tenant_id).order_by(Lead.created_date.desc())
    ).scalars().all()
    workers = get_workers()
    return render_template('business/leads.html', leads=all_leads, stages=LEAD_STAGES, workers=workers, today=date.today())


@business_bp.route('/business-hub/leads/add', methods=['POST'])
@handle_database_errors
def add_lead():
    r = require_login()
    if r: return r
    from app.models.lead import Lead
    try:
        lead = Lead(
            tenant_id=get_tenant_id(),
            name=sanitize_input(request.form.get('name', '')),
            phone=sanitize_input(request.form.get('phone', '')),
            email=sanitize_input(request.form.get('email', '')),
            source=sanitize_input(request.form.get('source', '')),
            stage=request.form.get('stage', 'Új lead'),
            notes=sanitize_input(request.form.get('notes', '')),
            assigned_worker_id=request.form.get('assigned_worker_id', type=int) or None,
            created_date=date.today(),
        )
        db.session.add(lead)
        db.session.commit()
        flash('Lead hozzáadva!', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    next_url = request.form.get('next', url_for('business.leads'))
    return redirect(next_url)


@business_bp.route('/business-hub/leads/<int:lead_id>/stage', methods=['POST'])
@handle_database_errors
def update_lead_stage(lead_id):
    r = require_login()
    if r: return r
    from app.models.lead import Lead
    lead = db.session.get(Lead, lead_id)
    if lead:
        lead.stage = request.form.get('stage', lead.stage)
        db.session.commit()
        flash('Lead státusz frissítve!', 'success')
    return redirect(url_for('business.leads'))


@business_bp.route('/business-hub/leads/<int:lead_id>/delete', methods=['POST'])
@handle_database_errors
def delete_lead(lead_id):
    r = require_login()
    if r: return r
    from app.models.lead import Lead
    lead = db.session.get(Lead, lead_id)
    if lead:
        db.session.delete(lead)
        db.session.commit()
        flash('Lead törölve!', 'success')
    return redirect(url_for('business.leads'))


# ─────────────────────────────────────────────────────────────────
# FELADATOK
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/tasks')
@handle_database_errors
def tasks():
    r = require_login()
    if r: return r
    from app.models.task import Task
    from sqlalchemy import and_
    tenant_id = get_tenant_id()
    all_tasks = db.session.execute(
        db.select(Task).where(Task.tenant_id == tenant_id).order_by(Task.done, Task.deadline)
    ).scalars().all()
    workers = get_workers()

    from app.models.job import Job
    jobs = db.session.execute(
        db.select(Job).where(Job.tenant_id == tenant_id).order_by(Job.job_id.desc()).limit(50)
    ).scalars().all()

    return render_template('business/tasks.html', tasks=all_tasks, workers=workers, jobs=jobs, today=date.today())


@business_bp.route('/business-hub/tasks/add', methods=['POST'])
@handle_database_errors
def add_task():
    r = require_login()
    if r: return r
    from app.models.task import Task
    try:
        deadline_str = request.form.get('deadline', '')
        t = Task(
            tenant_id=get_tenant_id(),
            title=sanitize_input(request.form.get('title', '')),
            description=sanitize_input(request.form.get('description', '')),
            assigned_worker_id=request.form.get('assigned_worker_id', type=int) or None,
            deadline=date.fromisoformat(deadline_str) if deadline_str else None,
            job_id=request.form.get('job_id', type=int) or None,
            done=False,
            created_date=date.today(),
        )
        db.session.add(t)
        db.session.commit()
        flash('Feladat hozzáadva!', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    next_url = request.form.get('next', url_for('business.tasks'))
    return redirect(next_url)


@business_bp.route('/business-hub/tasks/<int:task_id>/toggle', methods=['POST'])
@handle_database_errors
def toggle_task(task_id):
    r = require_login()
    if r: return r
    from app.models.task import Task
    t = db.session.get(Task, task_id)
    if t:
        t.done = not t.done
        db.session.commit()
    return redirect(url_for('business.tasks'))


@business_bp.route('/business-hub/tasks/<int:task_id>/delete', methods=['POST'])
@handle_database_errors
def delete_task(task_id):
    r = require_login()
    if r: return r
    from app.models.task import Task
    t = db.session.get(Task, task_id)
    if t:
        db.session.delete(t)
        db.session.commit()
        flash('Feladat törölve!', 'success')
    return redirect(url_for('business.tasks'))


# ─────────────────────────────────────────────────────────────────
# TELJESÍTÉSI IGAZOLÁSOK
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/confirmations')
@handle_database_errors
def confirmations():
    r = require_login()
    if r: return r
    from app.models.worker_payment import PerformanceConfirmation
    tenant_id = get_tenant_id()
    items = db.session.execute(
        db.select(PerformanceConfirmation).where(PerformanceConfirmation.tenant_id == tenant_id)
        .order_by(PerformanceConfirmation.month.desc())
    ).scalars().all()
    workers = get_workers()
    try:
        from app.models.job import Job
        from sqlalchemy import and_
        jobs = db.session.execute(
            db.select(Job).where(Job.tenant_id == tenant_id)
            .order_by(Job.job_id.desc()).limit(50)
        ).scalars().all()
    except Exception:
        jobs = []
    return render_template('business/confirmations.html', confirmations=items, workers=workers, jobs=jobs, today=date.today())


@business_bp.route('/business-hub/confirmations/add', methods=['POST'])
@handle_database_errors
def add_confirmation():
    r = require_login()
    if r: return r
    from app.models.worker_payment import PerformanceConfirmation
    import os
    from werkzeug.utils import secure_filename
    try:
        file = request.files.get('file')
        file_path = None
        if file and file.filename:
            upload_dir = os.path.join('app', 'static', 'uploads', 'confirmations')
            os.makedirs(upload_dir, exist_ok=True)
            filename = secure_filename(file.filename)
            full_path = os.path.join(upload_dir, filename)
            file.save(full_path)
            file_path = f'/static/uploads/confirmations/{filename}'
        c = PerformanceConfirmation(
            tenant_id=get_tenant_id(),
            worker_id=request.form.get('worker_id', type=int),
            month=request.form.get('month', date.today().strftime('%Y-%m')),
            file_path=file_path,
            status='Beadva',
            notes=sanitize_input(request.form.get('notes', '')),
            job_id=request.form.get('job_id', type=int) or None,
        )
        db.session.add(c)
        db.session.commit()
        flash('Teljesítési igazolás beadva!', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    return redirect(url_for('business.confirmations'))


@business_bp.route('/business-hub/confirmations/<int:conf_id>/approve', methods=['POST'])
@handle_database_errors
def approve_confirmation(conf_id):
    r = require_login()
    if r: return r
    from app.models.worker_payment import PerformanceConfirmation
    c = db.session.get(PerformanceConfirmation, conf_id)
    if c:
        c.status = 'Jóváhagyva'
        db.session.commit()
        flash('Teljesítési igazolás jóváhagyva!', 'success')
    return redirect(url_for('business.confirmations'))


# ─────────────────────────────────────────────────────────────────
# ÚJ MUNKÁS HOZZÁADÁSA
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/workers/add', methods=['POST'])
@handle_database_errors
def add_worker():
    r = require_login()
    if r: return r
    from app.models.user import User
    from app.models.tenant_membership import TenantMembership
    from werkzeug.security import generate_password_hash
    try:
        email = sanitize_input(request.form.get('email', ''))
        # Email egyediség ellenőrzés
        existing = db.session.execute(
            db.select(User).where(User.email == email)
        ).scalar_one_or_none()
        if existing:
            flash(f'Ez az email cím már foglalt: {email}', 'error')
            next_url = request.form.get('next', url_for('business.hub'))
            return redirect(next_url)

        password_raw = request.form.get('password', '')
        first_name = sanitize_input(request.form.get('first_name', ''))
        last_name = sanitize_input(request.form.get('last_name', ''))
        full_name = f"{first_name} {last_name}".strip() or email.split('@')[0]

        # Username egyediség biztosítása
        base_username = full_name
        username = base_username
        counter = 1
        while db.session.execute(db.select(User).where(User.username == username)).scalar_one_or_none():
            username = f"{base_username} {counter}"
            counter += 1

        user = User(
            username=username,
            email=email,
            password_hash=generate_password_hash(password_raw),
            is_active=True,
        )
        db.session.add(user)
        db.session.flush()  # user.user_id generálás

        membership = TenantMembership(
            tenant_id=get_tenant_id(),
            user_id=user.user_id,
            role=request.form.get('role', 'technician'),
            status='active',
        )
        db.session.add(membership)
        db.session.commit()
        flash(f'Munkás hozzáadva: {full_name}', 'success')
    except Exception as ex:
        db.session.rollback()
        flash(f'Hiba: {ex}', 'error')
    next_url = request.form.get('next', url_for('business.hub'))
    return redirect(next_url)

# ─────────────────────────────────────────────────────────────────
# TELJESÍTÉSI IGAZOLÁS PDF GENERÁLÁS
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/confirmations/<int:conf_id>/pdf')
@handle_database_errors
def confirmation_pdf(conf_id):
    """Teljesítési igazolás PDF letöltése - munkalap stílusban"""
    r = require_login()
    if r: return r

    from flask import make_response
    from io import BytesIO
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_RIGHT, TA_CENTER
    import datetime, os

    from app.models.worker_payment import PerformanceConfirmation
    from app.models.tenant import Tenant

    c = db.session.get(PerformanceConfirmation, conf_id)
    if not c:
        flash('Igazolás nem található', 'error')
        return redirect(url_for('business.confirmations'))

    tenant_id = get_tenant_id()
    tenant = Tenant.find_by_id(tenant_id)
    settings = tenant.settings or {} if tenant else {}

    # Font betöltés — ugyanolyan mint technician.py-ban
    unicode_font = 'Helvetica'
    unicode_font_bold = 'Helvetica-Bold'
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        font_search = [
            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
            '/usr/share/fonts/truetype/freefont/FreeSans.ttf',
        ]
        bold_search = [
            '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
            '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
            '/usr/share/fonts/truetype/freefont/FreeSansBold.ttf',
        ]
        for fp, fb in zip(font_search, bold_search):
            if os.path.exists(fp) and os.path.exists(fb):
                try:
                    pdfmetrics.registerFont(TTFont('UniFont', fp))
                    pdfmetrics.registerFont(TTFont('UniFont-Bold', fb))
                    unicode_font = 'UniFont'
                    unicode_font_bold = 'UniFont-Bold'
                    break
                except Exception:
                    continue
    except Exception:
        pass

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                           rightMargin=1.5*cm, leftMargin=1.5*cm,
                           topMargin=1.5*cm, bottomMargin=1.5*cm)
    styles = getSampleStyleSheet()
    for style in styles.byName.values():
        style.fontName = unicode_font

    story = []
    white = colors.white
    primary = colors.HexColor('#1e3a5f')
    accent  = colors.HexColor('#e87e04')

    # ── FEJLÉC ────────────────────────────────────────────────────
    doc_num = f"TI-{datetime.date.today().year}-{conf_id:04d}"
    company_name = tenant.name if tenant else 'STAR LABS Kft.'

    header_data = [[
        Paragraph(f'<font color="white" size="20"><b>{company_name}</b></font>', styles['Normal']),
        Paragraph(f'<font color="white" size="11"><b>TELJESÍTÉSI IGAZOLÁS</b><br/>{doc_num}<br/>{datetime.date.today().strftime("%Y. %m. %d.")}</font>', styles['Normal'])
    ]]
    header_table = Table(header_data, colWidths=[10*cm, 8*cm])
    header_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), accent),
        ('TOPPADDING', (0,0), (-1,-1), 16), ('BOTTOMPADDING', (0,0), (-1,-1), 16),
        ('LEFTPADDING', (0,0), (-1,-1), 16), ('RIGHTPADDING', (0,0), (-1,-1), 16),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'), ('ALIGN', (1,0), (1,0), 'RIGHT'),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 0.5*cm))

    # ── IGAZOLÁS ADATOK ───────────────────────────────────────────
    info_style = ParagraphStyle('info', fontSize=9, leading=14, fontName=unicode_font)

    worker_name = c.worker.username if c.worker else '—'
    job_info = f'#{c.job_id}' if c.job_id else '—'

    seller_addr = tenant.address if tenant else ''
    tax_id = settings.get('tax_id', '')

    left_text = f'<b>Kibocsátó</b><br/>{company_name}<br/>{seller_addr}'
    if tax_id:
        left_text += f'<br/>Adószám: {tax_id}'

    right_text = f'<b>Munkás</b><br/>{worker_name}<br/><b>Hónap:</b> {c.month}<br/><b>Kapcsolódó munka:</b> {job_info}'

    info_data = [[Paragraph(left_text, info_style), Paragraph(right_text, info_style)]]
    info_table = Table(info_data, colWidths=[9*cm, 9*cm])
    info_table.setStyle(TableStyle([
        ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('LINEAFTER', (0,0), (0,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('PADDING', (0,0), (-1,-1), 12), ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(info_table)
    story.append(Spacer(1, 0.5*cm))

    # ── IGAZOLÁS TARTALOM ─────────────────────────────────────────
    th_style = ParagraphStyle('th', fontSize=9, textColor=white, fontName=unicode_font_bold)
    cell_style = ParagraphStyle('cell', fontSize=9, fontName=unicode_font)

    rows = [[
        Paragraph('Megnevezés', th_style),
        Paragraph('Érték', th_style),
    ]]

    status_label = '✓ Jóváhagyva' if c.status == 'Jóváhagyva' else '⏳ Beadva / Feldolgozás alatt'
    rows.append([Paragraph('Státusz', cell_style), Paragraph(status_label, cell_style)])
    rows.append([Paragraph('Hónap', cell_style), Paragraph(c.month, cell_style)])
    rows.append([Paragraph('Munkás', cell_style), Paragraph(worker_name, cell_style)])
    rows.append([Paragraph('Dokumentum száma', cell_style), Paragraph(doc_num, cell_style)])

    if c.job_id:
        # Ha van kapcsolódó munka, adjuk hozzá a részleteket
        try:
            from app.models.job import Job
            job = db.session.get(Job, c.job_id)
            if job:
                customer = job.customer_rel
                buyer_name = f"{customer.first_name} {customer.family_name}" if customer else '—'
                rows.append([Paragraph('Ügyfél', cell_style), Paragraph(buyer_name, cell_style)])
                rows.append([Paragraph('Munka dátuma', cell_style), Paragraph(job.job_date.strftime('%Y. %m. %d.') if job.job_date else '—', cell_style)])
                total = f"{int(job.total_cost):,} Ft".replace(',', ' ') if job.total_cost else '—'
                rows.append([Paragraph('Munka értéke', cell_style), Paragraph(total, cell_style)])

                # Tételek
                services = job.get_services()
                parts = job.get_parts()
                if services or parts:
                    rows.append([Paragraph('<b>Elvégzett munkák:</b>', cell_style), Paragraph('', cell_style)])
                    for s in services:
                        rows.append([Paragraph(f'  • {s["service_name"]}', cell_style),
                                     Paragraph(f'{s["qty"]} db × {int(s["cost"]):,} Ft'.replace(',', ' '), cell_style)])
                    for p in parts:
                        rows.append([Paragraph(f'  • {p["part_name"]}', cell_style),
                                     Paragraph(f'{p["qty"]} db × {int(p["cost"]):,} Ft'.replace(',', ' '), cell_style)])
        except Exception:
            pass

    if c.notes:
        rows.append([Paragraph('Megjegyzés', cell_style), Paragraph(c.notes, cell_style)])

    t = Table(rows, colWidths=[6*cm, 12*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), primary),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [white, colors.HexColor('#f8fafc')]),
        ('GRID', (0,0), (-1,-1), 0.3, colors.HexColor('#e2e8f0')),
        ('PADDING', (0,0), (-1,-1), 8),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ]))
    story.append(t)
    story.append(Spacer(1, 1*cm))

    # ── ALÁÍRÁS BLOKK ─────────────────────────────────────────────
    sig_style = ParagraphStyle('sig', fontSize=9, fontName=unicode_font, alignment=TA_CENTER)
    sig_data = [[
        Paragraph('____________________________<br/><br/>Kibocsátó aláírása', sig_style),
        Paragraph('____________________________<br/><br/>Munkás aláírása', sig_style),
        Paragraph('____________________________<br/><br/>Jóváhagyó aláírása', sig_style),
    ]]
    sig_table = Table(sig_data, colWidths=[6*cm, 6*cm, 6*cm])
    sig_table.setStyle(TableStyle([
        ('PADDING', (0,0), (-1,-1), 12),
        ('VALIGN', (0,0), (-1,-1), 'BOTTOM'),
    ]))
    story.append(sig_table)
    story.append(Spacer(1, 0.5*cm))

    # ── LÁBLÉC ────────────────────────────────────────────────────
    story.append(HRFlowable(width='100%', thickness=0.5, color=colors.HexColor('#e2e8f0')))
    footer_style = ParagraphStyle('footer', fontSize=8, textColor=colors.HexColor('#94a3b8'),
                                   alignment=TA_CENTER, fontName=unicode_font)
    story.append(Paragraph(f'Powered by RepairOS · {company_name}', footer_style))

    doc.build(story)
    buffer.seek(0)

    response = make_response(buffer.read())
    response.headers['Content-Type'] = 'application/pdf'
    response.headers['Content-Disposition'] = f'inline; filename=teljesites-igazolas-{doc_num}.pdf'
    return response


# ─────────────────────────────────────────────────────────────────
# AI AJÁNLATKÉSZÍTÉS
# ─────────────────────────────────────────────────────────────────

@business_bp.route('/business-hub/generate-quote', methods=['POST'])
@handle_database_errors
def generate_quote():
    """AI-alapú ajánlat PPTX generálás"""
    r = require_login()
    if r: return r

    import os, base64, json, requests as req_lib
    from flask import make_response
    from io import BytesIO

    tenant_id = get_tenant_id()

    # ── FORM ADATOK ───────────────────────────────────────────────
    quote_title  = sanitize_input(request.form.get('quote_title', 'Ajánlat'))
    customer_id  = request.form.get('customer_id', type=int)
    valid_until  = request.form.get('valid_until', '')
    ai_notes     = sanitize_input(request.form.get('ai_notes', ''))
    item_names   = request.form.getlist('item_name[]')
    item_qtys    = request.form.getlist('item_qty[]')
    item_prices  = request.form.getlist('item_price[]')

    items = []
    total = 0
    for name, qty, price in zip(item_names, item_qtys, item_prices):
        if name:
            q = int(qty or 1)
            p = float(price or 0)
            subtotal = q * p
            total += subtotal
            items.append({'name': name, 'qty': q, 'price': p, 'subtotal': subtotal})

    # ── ÜGYFÉL ADATOK ─────────────────────────────────────────────
    customer_name = 'Tisztelt Ügyfél'
    try:
        from app.models.customer import Customer
        from app.extensions import db
        if customer_id:
            c = db.session.get(Customer, customer_id)
            if c:
                customer_name = f"{c.first_name} {c.family_name}"
    except Exception:
        pass

    # ── VÁLLALAT ADATOK ───────────────────────────────────────────
    try:
        from app.models.tenant import Tenant
        tenant = Tenant.find_by_id(tenant_id)
        company_name = tenant.name if tenant else 'STAR LABS Kft.'
        company_addr = tenant.address if tenant else 'Budapest'
        settings = tenant.settings or {} if tenant else {}
        company_tax = settings.get('tax_id', '')
        company_bank = settings.get('bank_account', '')
        company_email = tenant.email if tenant else ''
        company_phone = tenant.phone if tenant else ''
    except Exception:
        company_name = 'STAR LABS Kft.'
        company_addr = 'Budapest'
        company_tax = company_bank = company_email = company_phone = ''

    # ── KÉPEK BEOLVASÁSA ──────────────────────────────────────────
    logo_b64 = None
    logo_file = request.files.get('logo')
    if logo_file and logo_file.filename:
        logo_b64 = base64.b64encode(logo_file.read()).decode('utf-8')
        logo_mime = 'image/png' if logo_file.filename.lower().endswith('.png') else 'image/jpeg'

    image_files = request.files.getlist('images')
    images_b64 = []
    for img in image_files[:3]:
        if img and img.filename:
            images_b64.append({
                'data': base64.b64encode(img.read()).decode('utf-8'),
                'mime': 'image/png' if img.filename.lower().endswith('.png') else 'image/jpeg',
                'name': img.filename
            })

    # ── ANTHROPIC API HÍVÁS ───────────────────────────────────────
    anthropic_key = os.environ.get('ANTHROPIC_API_KEY', '')
    ai_texts = {'intro': '', 'items': [], 'closing': '', 'payment_schedule': [], '_detailed': False}

    if anthropic_key:
        items_text = '\n'.join([f"- {i['name']}: {i['qty']} db × {int(i['price']):,} Ft = {int(i['subtotal']):,} Ft" for i in items]) if items else 'Nincs megadva'
        total_str = f"{int(total):,} Ft" if total else 'Lásd részletek'
        detailed_mode = len(ai_notes) > 100
        logger.error(f"AI quote DEBUG: notes_len={len(ai_notes)}, detailed={detailed_mode}, items={len(items)}, has_key={bool(anthropic_key)}")

        if detailed_mode:
            prompt = f"""Te egy profi magyar üzleti ajánlatszöveg-író vagy. Készíts PPTX prezentáció tartalmat az alábbi RÉSZLETES INSTRUKCIÓK alapján - ezeket KÖTELEZŐ követni!

Cég: {company_name}
Ajánlat címe: {quote_title}
Ügyfél: {customer_name}
Tételek: {items_text}
Teljes összeg: {total_str}
Érvényesség: {valid_until}

RÉSZLETES INSTRUKCIÓK (KÖTELEZŐ KÖVETNI):
{ai_notes}

Válaszolj CSAK JSON formátumban, semmi más:
{{
  "intro": "2-3 mondatos bevezető (megszólítással, üzleti hangvételű)",
  "project_description": "2-3 mondatos projekt leírás",
  "technical_content": "3-4 mondatos műszaki tartalom",
  "payment_schedule": [
    {{"milestone": "mérföldkő neve", "percent": 15, "amount": "15 000 000 Ft", "description": "rövid leírás"}}
  ],
  "timeline": "teljesítési határidő leírása",
  "warranty": "garanciális feltételek",
  "payment_terms": "fizetési feltételek",
  "legal_terms": "általános jogi feltételek",
  "closing": "2-3 mondatos záró szöveg"
}}"""
        else:
            prompt = f"""Te egy profi üzleti ajánlatszöveg-írói asszisztens vagy. Írj rövid, meggyőző magyar nyelvű szövegeket egy PPTX ajánlathoz.

Cég: {company_name}
Ajánlat címe: {quote_title}
Ügyfél: {customer_name}
Tételek: {items_text}
Összesen: {total_str}
Érvényes: {valid_until}
{('Megjegyzés: ' + ai_notes) if ai_notes else ''}

Válaszolj CSAK JSON formátumban, semmi más:
{{
  "intro": "2-3 mondatos bevezető szöveg az ajánlathoz (megszólítással)",
  "items": [{{"name": "tétel neve", "description": "1 mondatos leírás"}}],
  "closing": "2-3 mondatos záró szöveg"
}}"""

        try:
            api_resp = req_lib.post(
                'https://api.anthropic.com/v1/messages',
                headers={
                    'x-api-key': anthropic_key,
                    'anthropic-version': '2023-06-01',
                    'content-type': 'application/json',
                },
                json={
                    'model': 'claude-sonnet-4-20250514',
                    'max_tokens': 2000,
                    'messages': [{'role': 'user', 'content': prompt}]
                },
                timeout=45
            )
            if api_resp.status_code == 200:
                raw = api_resp.json()['content'][0]['text'].strip()
                logger.error(f"AI API OK: raw[:100]={raw[:100]}")
                if raw.startswith('```'):
                    raw = raw.split('```')[1]
                    if raw.startswith('json'): raw = raw[4:]
                parsed = json.loads(raw.strip())
                ai_texts.update(parsed)
                ai_texts['_detailed'] = detailed_mode
                logger.error(f"AI texts keys: {list(ai_texts.keys())}, _detailed={ai_texts.get('_detailed')}")
            else:
                logger.error(f"AI API hiba: {api_resp.status_code} {api_resp.text[:200]}")
        except Exception as e:
            logger.error(f"AI szöveggenerálás hiba: {e}")

    # ── PPTX GENERÁLÁS python-pptx-szel ──────────────────────────
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import datetime as dt

        DARK   = RGBColor(0x1e, 0x3a, 0x5f)
        ACCENT = RGBColor(0xe8, 0x7e, 0x04)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY   = RGBColor(0x64, 0x74, 0x8b)
        LIGHT  = RGBColor(0xf8, 0xfa, 0xfc)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]  # blank

        def add_rect(slide, x, y, w, h, color):
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            return shape

        def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            return txBox

        def add_image_b64(slide, b64data, mime, x, y, w, h):
            try:
                img_bytes = base64.b64decode(b64data)
                img_stream = BytesIO(img_bytes)
                slide.shapes.add_picture(img_stream, Inches(x), Inches(y), Inches(w), Inches(h))
            except Exception as ie:
                logger.warning(f"Kép beillesztés hiba: {ie}")

        # ════ SLIDE 1: BORÍTÓ ════════════════════════════════════
        sl1 = prs.slides.add_slide(blank)
        add_rect(sl1, 0, 0, 13.33, 7.5, DARK)
        add_rect(sl1, 0, 5.8, 13.33, 1.7, ACCENT)

        # Logo ha van
        if logo_b64:
            add_image_b64(sl1, logo_b64, logo_mime, 0.5, 0.3, 2.5, 1.0)

        add_text(sl1, company_name, 0.5, 1.5, 12, 0.8, size=16, color=RGBColor(0x93,0xc5,0xfd))
        add_text(sl1, quote_title, 0.5, 2.3, 12, 1.5, size=40, bold=True, color=WHITE)
        add_text(sl1, f'Ajánlat: {customer_name}', 0.5, 4.0, 8, 0.6, size=20, color=RGBColor(0xfb,0xd3,0x8d))

        if valid_until:
            add_text(sl1, f'Érvényes: {valid_until}', 0.5, 4.7, 6, 0.5, size=14, color=GRAY)

        add_text(sl1, dt.date.today().strftime('%Y. %m. %d.'), 10, 4.7, 2.5, 0.5, size=13, color=GRAY, align=PP_ALIGN.RIGHT)

        # Képek a borítón ha van
        if images_b64:
            add_image_b64(sl1, images_b64[0]['data'], images_b64[0]['mime'], 9.5, 0.5, 3.3, 5.0)

        # ════ SLIDE 2: BEVEZETŐ ══════════════════════════════════
        sl2 = prs.slides.add_slide(blank)
        add_rect(sl2, 0, 0, 13.33, 1.2, DARK)
        add_text(sl2, 'Bemutatkozás & Ajánlatunk', 0.4, 0.2, 12, 0.8, size=28, bold=True)

        intro_text = ai_texts.get('intro') or f'Tisztelt {customer_name}! Örömmel küldjük Önnek az alábbi árajánlatot. Bízunk benne, hogy megoldásaink megfelelnek elvárásainak.'
        add_text(sl2, intro_text, 0.5, 1.5, 8, 3.0, size=16, color=RGBColor(0x1e,0x29,0x3b))

        # Stat kártyák
        add_rect(sl2, 0.5, 5.0, 3.5, 1.8, LIGHT)
        add_text(sl2, str(len(items)), 0.7, 5.1, 3.0, 0.7, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(sl2, 'Ajánlott tétel', 0.7, 5.8, 3.0, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

        add_rect(sl2, 4.5, 5.0, 4.0, 1.8, LIGHT)
        add_text(sl2, f'{int(total):,} Ft'.replace(',', ' '), 4.7, 5.1, 3.5, 0.7, size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(sl2, 'Nettó összérték', 4.7, 5.8, 3.5, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

        if images_b64 and len(images_b64) > 0:
            add_image_b64(sl2, images_b64[0]['data'], images_b64[0]['mime'], 9.0, 1.5, 3.8, 3.0)

        # ════ SLIDE 3: TÉTELEK ═══════════════════════════════════
        sl3 = prs.slides.add_slide(blank)
        add_rect(sl3, 0, 0, 13.33, 1.2, ACCENT)
        add_text(sl3, 'Ajánlat tételei', 0.4, 0.2, 12, 0.8, size=28, bold=True)

        # Táblázat fejléc
        add_rect(sl3, 0.3, 1.3, 7.5, 0.45, DARK)
        add_text(sl3, 'Megnevezés', 0.35, 1.32, 4.5, 0.4, size=12, bold=True)
        add_text(sl3, 'Menny.', 4.85, 1.32, 1.0, 0.4, size=12, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl3, 'Egységár', 5.85, 1.32, 1.2, 0.4, size=12, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl3, 'Részösszeg', 7.05, 1.32, 1.5, 0.4, size=12, bold=True, align=PP_ALIGN.RIGHT)

        ai_item_map = {i['name']: i.get('description','') for i in (ai_texts.get('items') or [])}
        row_h = 0.55
        for idx, item in enumerate(items[:8]):
            y = 1.85 + idx * row_h
            bg = LIGHT if idx % 2 == 0 else WHITE
            add_rect(sl3, 0.3, y, 7.5, row_h - 0.05, bg)
            desc = ai_item_map.get(item['name'], '')
            name_text = item['name'] + (f'\n  {desc}' if desc else '')
            add_text(sl3, name_text, 0.35, y+0.03, 4.4, row_h-0.08, size=11, color=RGBColor(0x1e,0x29,0x3b))
            add_text(sl3, f"{item['qty']} db", 4.85, y+0.05, 1.0, 0.4, size=11, color=DARK, align=PP_ALIGN.CENTER)
            add_text(sl3, f"{int(item['price']):,} Ft".replace(',', ' '), 5.85, y+0.05, 1.2, 0.4, size=11, color=DARK, align=PP_ALIGN.CENTER)
            add_text(sl3, f"{int(item['subtotal']):,} Ft".replace(',', ' '), 7.05, y+0.05, 1.5, 0.4, size=11, color=DARK, align=PP_ALIGN.RIGHT)

        # Összesen sor
        total_y = 1.85 + min(len(items), 8) * row_h + 0.1
        add_rect(sl3, 0.3, total_y, 7.5, 0.55, DARK)
        add_text(sl3, 'ÖSSZESEN (nettó)', 0.35, total_y+0.08, 5.5, 0.4, size=13, bold=True)
        add_text(sl3, f"{int(total):,} Ft".replace(',', ' '), 6.0, total_y+0.08, 1.5, 0.4, size=13, bold=True, align=PP_ALIGN.RIGHT)

        if images_b64 and len(images_b64) > 1:
            add_image_b64(sl3, images_b64[1]['data'], images_b64[1]['mime'], 8.5, 1.5, 4.3, 5.5)

        # ════ SLIDE 4: ZÁRÁS / KONTAKT ═══════════════════════════
        sl4 = prs.slides.add_slide(blank)
        add_rect(sl4, 0, 0, 13.33, 7.5, DARK)
        add_rect(sl4, 0, 5.5, 13.33, 2.0, ACCENT)

        closing = ai_texts.get('closing') or 'Kérjük, forduljon hozzánk bizalommal! Csapatunk készen áll az ajánlat részletes megbeszélésére.'
        add_text(sl4, 'Következő lépések', 0.5, 0.5, 12, 0.8, size=30, bold=True)
        add_text(sl4, closing, 0.5, 1.5, 9, 2.0, size=16, color=RGBColor(0xcb,0xd5,0xe1))

        # Kontakt adatok
        contact = company_name
        if company_addr:   contact += f'\n{company_addr}'
        if company_email:  contact += f'\n{company_email}'
        if company_phone:  contact += f'\n{company_phone}'
        if company_tax:    contact += f'\nAdószám: {company_tax}'
        add_text(sl4, contact, 0.5, 5.55, 8, 1.8, size=13, color=WHITE)
        add_text(sl4, f'Érvényes: {valid_until}' if valid_until else '', 9.5, 5.7, 3.3, 0.5, size=13, color=WHITE, align=PP_ALIGN.RIGHT)

        if logo_b64:
            add_image_b64(sl4, logo_b64, logo_mime, 10.5, 0.5, 2.3, 0.9)

        if images_b64 and len(images_b64) > 2:
            add_image_b64(sl4, images_b64[2]['data'], images_b64[2]['mime'], 9.0, 1.5, 3.8, 3.5)

        # ════ RÉSZLETES MÓD: extra slide-ok ══════════════════════
        if ai_texts.get('_detailed'):

            # ── SLIDE 5: PROJEKT LEÍRÁS + MŰSZAKI TARTALOM ──────
            sl5 = prs.slides.add_slide(blank)
            add_rect(sl5, 0, 0, 13.33, 1.2, DARK)
            add_text(sl5, 'Projekt leírás & Műszaki tartalom', 0.4, 0.2, 12, 0.8, size=26, bold=True)

            proj_desc = ai_texts.get('project_description', '')
            tech = ai_texts.get('technical_content', '')

            add_rect(sl5, 0.3, 1.4, 6.0, 0.4, ACCENT)
            add_text(sl5, 'Projekt leírás', 0.4, 1.42, 5.5, 0.35, size=13, bold=True)
            add_text(sl5, proj_desc, 0.4, 1.9, 6.0, 2.5, size=13, color=RGBColor(0x1e,0x29,0x3b))

            add_rect(sl5, 0.3, 4.5, 6.0, 0.4, DARK)
            add_text(sl5, 'Műszaki tartalom', 0.4, 4.52, 5.5, 0.35, size=13, bold=True)
            add_text(sl5, tech, 0.4, 5.0, 6.0, 2.0, size=13, color=RGBColor(0x1e,0x29,0x3b))

            if images_b64 and len(images_b64) > 0:
                add_image_b64(sl5, images_b64[0]['data'], images_b64[0]['mime'], 7.0, 1.4, 5.8, 5.5)

            # ── SLIDE 6: FIZETÉSI ÜTEMEZÉS ───────────────────────
            payment_sched = ai_texts.get('payment_schedule', [])
            if payment_sched:
                sl6 = prs.slides.add_slide(blank)
                add_rect(sl6, 0, 0, 13.33, 1.2, ACCENT)
                add_text(sl6, 'Fizetési ütemezés', 0.4, 0.2, 12, 0.8, size=28, bold=True)

                # Fejléc
                add_rect(sl6, 0.3, 1.3, 12.5, 0.45, DARK)
                add_text(sl6, 'Mérföldkő', 0.4, 1.32, 5.5, 0.4, size=12, bold=True)
                add_text(sl6, '%', 5.95, 1.32, 1.0, 0.4, size=12, bold=True, align=PP_ALIGN.CENTER)
                add_text(sl6, 'Összeg', 7.0, 1.32, 2.5, 0.4, size=12, bold=True, align=PP_ALIGN.CENTER)
                add_text(sl6, 'Leírás', 9.6, 1.32, 3.0, 0.4, size=12, bold=True)

                for idx, ps in enumerate(payment_sched[:8]):
                    y = 1.85 + idx * 0.6
                    bg = LIGHT if idx % 2 == 0 else WHITE
                    add_rect(sl6, 0.3, y, 12.5, 0.55, bg)
                    add_text(sl6, ps.get('milestone',''), 0.4, y+0.08, 5.5, 0.4, size=11, color=DARK, bold=True)
                    add_text(sl6, f"{ps.get('percent',0)}%", 5.95, y+0.08, 1.0, 0.4, size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
                    add_text(sl6, ps.get('amount',''), 7.0, y+0.08, 2.5, 0.4, size=11, color=DARK, align=PP_ALIGN.CENTER)
                    add_text(sl6, ps.get('description',''), 9.6, y+0.08, 3.0, 0.4, size=10, color=GRAY)

                # Összesen
                tot_y = 1.85 + min(len(payment_sched), 8) * 0.6 + 0.1
                add_rect(sl6, 0.3, tot_y, 12.5, 0.55, DARK)
                add_text(sl6, 'ÖSSZESEN', 0.4, tot_y+0.08, 5.5, 0.4, size=13, bold=True)
                add_text(sl6, '100%', 5.95, tot_y+0.08, 1.0, 0.4, size=13, bold=True, align=PP_ALIGN.CENTER)
                total_disp = f"{int(total):,} Ft".replace(',', ' ') if total else ai_texts.get('payment_schedule', [{}])[-1].get('amount','')
                add_text(sl6, total_disp, 7.0, tot_y+0.08, 2.5, 0.4, size=13, bold=True, align=PP_ALIGN.CENTER)

            # ── SLIDE 7: FELTÉTELEK ──────────────────────────────
            sl7 = prs.slides.add_slide(blank)
            add_rect(sl7, 0, 0, 13.33, 1.2, DARK)
            add_text(sl7, 'Szerződési feltételek', 0.4, 0.2, 12, 0.8, size=26, bold=True)

            conditions = [
                ('Teljesítési határidő', ai_texts.get('timeline', '')),
                ('Garanciális feltételek', ai_texts.get('warranty', '')),
                ('Fizetési feltételek', ai_texts.get('payment_terms', '')),
                ('Általános feltételek', ai_texts.get('legal_terms', '')),
            ]
            y = 1.4
            for label, text in conditions:
                if text:
                    add_rect(sl7, 0.3, y, 12.5, 0.35, ACCENT)
                    add_text(sl7, label, 0.4, y+0.03, 12.0, 0.3, size=12, bold=True)
                    y += 0.4
                    add_text(sl7, text, 0.4, y, 12.0, 0.7, size=12, color=RGBColor(0x1e,0x29,0x3b))
                    y += 0.85

        # ── MENTÉS ────────────────────────────────────────────────
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Ékezetek eltávolítása a fájlnévből
        import unicodedata
        safe_title = unicodedata.normalize('NFKD', quote_title)
        safe_title = ''.join(c for c in safe_title if not unicodedata.combining(c))
        safe_title = ''.join(c if c.isalnum() or c in '-_ ' else '_' for c in safe_title)
        safe_title = safe_title.replace(' ', '_')[:30].strip('_')

        response = make_response(pptx_buffer.read())
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        response.headers['Content-Disposition'] = f'attachment; filename=ajanlat_{safe_title}.pptx'
        return response

    except ImportError:
        flash('python-pptx csomag hiányzik! pip install python-pptx', 'error')
        return redirect(url_for('business.hub'))
    except Exception as e:
        logger.error(f"PPTX generálás hiba: {e}", exc_info=True)
        flash(f'PPTX generálás hiba: {e}', 'error')
        return redirect(url_for('business.hub'))
